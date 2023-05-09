# -*- coding: utf-8 -*-
from odoo import http
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Pt
from datetime import datetime, date, time, timedelta
import calendar

class Powerpoint(http.Controller):
    @http.route('/ppt/<id>', type='http', auth='public')
    def index(self,id):
        partida = http.request.env['partidas.partidas']
        orden = partida.sudo().search([('id','=',id)])
        prs = Presentation('/usr/lib/python3/dist-packages/odoo/odoo-extra-addons/powerpoint/static/plantilla.pptx')


        ahora = datetime.now()
        months = ["Unknown",
                "Enero",
                "Febrero",
                "Marzo",
                "Abril",
                "Mayo",
                "Junio",
                "Julio",
                "Agosto",
                "Septiembre",
                "Octubre",
                "Noviembre",
                "Diciembre"]
        fecha = str(ahora.day) + " de " + months[ahora.month] + " de "+ str(ahora.year)
        comentarios = ""
        for i in orden.avance_semaforo:
            comentarios = i.comentarios_generales

        valores = {
                    'titulo': orden.obra.descripcion,
                    'fecha': fecha,
                    'localidad': orden.localidad,
                    'municipio': orden.municipio.municipio_delegacion,
                    'contrato': orden.nombre_partida,
                    'monto': orden.total_partida,
                    'avfis': orden.a_fis,
                    'avprog': orden.porcentajeProgramado,
                    'atraso': orden.atraso,
                    'comentarios': comentarios,
                }

        for slide in prs.slides:
            for shape in slide.shapes:
                for key in valores:
                    if '%{}%'.format(key) in shape.text:
                        new_paragraph = shape.text.replace('%{}%'.format(key),str(valores[key]))
                        shape.text = new_paragraph
                        #shape.text_frame.fit_text()
                        for d in shape.text_frame.paragraphs:
                            d.font.size = Pt(10)
                        #shape.text_frame.paragraphs[0].font.size = Pt(12)

        prs.save('/tmp/test.pptx')
        f = open('/tmp/test.pptx', mode="rb")
        return http.request.make_response(f.read(),
                                          [('Content-Type', 'application/octet-stream'),
                                           ('Content-Disposition',
                                            'attachment; filename="{}"'.format('sidur.pptx'))
                                           ])

    # @http.route('/ppt/', auth='public')
    # def index(self, **kw):
    #     estimaciones = http.request.env['control.estimaciones']
    #     orden = estimaciones.sudo().search([('obra','=',kw['id'])])

    #     resultado = []
    #     obra = ''
    #     for i in orden:
    #         resultado.append(i.id)
    #         obra = i.obra.numero_contrato.contrato
    #         print(obra)

    #     # df1 = pd.DataFrame({'ID Estimación': resultado})
    #     df1 = pd.DataFrame({'ID Estimación': resultado, 'Obra': obra, 'C': 'fg', 'D': 'sdf', 'F': 'fgfg'})

    #     # Create a Pandas Excel writer using XlsxWriter as the engine.
    #     writer = pd.ExcelWriter('/tmp/estimaciones.xlsx', engine='xlsxwriter')

    #     # Write each dataframe to a different worksheet.
    #     df1.to_excel(writer, sheet_name='Sheet1')

    #     # Close the Pandas Excel writer and output the Excel file.
    #     writer.save()

    #     f = open('/tmp/estimaciones.xlsx', mode="rb")
    #     return http.request.make_response(f.read(),
    #                                       [('Content-Type', 'application/octet-stream'),
    #                                        ('Content-Disposition',
    #                                         'attachment; filename="{}"'.format('estimaciones.xls'))
    #                                        ])